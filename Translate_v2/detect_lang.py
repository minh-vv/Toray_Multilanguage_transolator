import os
import docx
import openpyxl
import PyPDF2
import pptx
from fastapi import HTTPException
from google.cloud import translate
from dotenv import load_dotenv

load_dotenv()
# Định nghĩa các ngôn ngữ hỗ trợ
VIETNAMESE_SYMBOL = 'vi'
JAPANESE_SYMBOL = 'ja'
CHINESE_SYMBOL = 'zh-CN'
ENGLISH_SYMBOL = 'en'

# Lấy giá trị biến môi trường
google_credentials_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
PROJECT_ID = os.getenv("PROJECT_ID")
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = google_credentials_path

# Định nghĩa các ngôn ngữ hỗ trợ
LANGUAGES = {
    'vi': "Vietnamese",
    'ja': "Japanese",
    # 'zh-CN': "Chinese",
    # 'en': "English"
}
language_pair = {
    "vi-en": 1,
    "vi-ja": 2,
    "vi-zh-CN": 3,
    "en-ja": 4,
    "en-zh-CN": 5,
    "ja-zh-CN": 6
}
def extract_content(file_path: str) -> str:
    """Trích xuất nội dung từ file (PDF, DOCX, XLSX) để phát hiện ngôn ngữ"""
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    file_extension = file_path.split(".")[-1].lower()

    if file_extension == "pdf":
        return extract_from_pdf(file_path)
    elif file_extension == "docx":
        return extract_from_docx(file_path)
    elif file_extension == "xlsx":
        return extract_from_xlsx(file_path)
    elif file_extension == "pptx":
        return extract_from_pptx(file_path)
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")
def extract_from_pptx(file_path: str) -> str:
    """Trích xuất văn bản từ file PPTX"""
    content = ""
    prs = pptx.Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + " "
                if len(content.split()) >= 100:
                    break
        if len(content.split()) >= 100:
            break
    return get_first_100_words(content)
def get_first_100_words(text: str) -> str:
    """Lấy 100 từ đầu tiên của văn bản"""
    words = text.split()[:100]
    return ' '.join(words)

def extract_from_pdf(file_path: str) -> str:
    """Trích xuất văn bản từ file PDF"""
    content = ""
    with open(file_path, "rb") as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page in pdf_reader.pages:
            content += page.extract_text() or ""
            if len(content.split()) >= 100:
                break
    return get_first_100_words(content)

def extract_from_docx(file_path: str) -> str:
    """Trích xuất văn bản từ file DOCX"""
    content = ""
    doc = docx.Document(file_path)
    for para in doc.paragraphs:
        content += para.text + " "
        if len(content.split()) >= 100:
            break
    return get_first_100_words(content)

def extract_from_xlsx(file_path: str) -> str:
    """Trích xuất văn bản từ file XLSX"""
    content = ""
    workbook = openpyxl.load_workbook(file_path)
    sheet = workbook[workbook.sheetnames[0]]

    for row in sheet.iter_rows(values_only=True):
        for cell in row:
            if cell:
                content += str(cell) + " "
                if len(content.split()) >= 100:
                    break
        if len(content.split()) >= 100:
            break

    return get_first_100_words(content)

def detect_language(text: str):
    """Phát hiện ngôn ngữ của file"""
    try:
        client = translate.TranslationServiceClient()
        parent = f"projects/{PROJECT_ID}/locations/us-central1"
        request = translate.DetectLanguageRequest(content=text, parent=parent)

        response = client.detect_language(request=request)
        detected_language = response.languages[0].language_code
        return detected_language
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error detecting language: {str(e)}")
